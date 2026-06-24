from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Brand colours ──────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy background
CREAM  = RGBColor(0xF5, 0xF0, 0xE8)   # warm cream text
GOLD   = RGBColor(0xE8, 0xC4, 0x6A)   # gold accent
TEAL   = RGBColor(0x2D, 0x9C, 0x9C)   # teal accent
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xD0, 0xD8, 0xE8)   # light blue-grey for body text

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── helpers ────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb=None, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=24, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def bg(slide, color=NAVY):
    """Fill the entire slide background."""
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=color)


def accent_bar(slide, color=GOLD, x=Inches(0), y=Inches(6.9),
               w=Inches(13.33), h=Inches(0.08)):
    add_rect(slide, x, y, w, h, fill_rgb=color)


def slide_number(slide, n, total=10):
    add_text(slide, f"{n} / {total}", Inches(12.2), Inches(7.1),
             Inches(1), Inches(0.35), font_size=10, color=LIGHT,
             align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 – TITLE
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

# left teal panel
add_rect(s, 0, 0, Inches(5.2), SLIDE_H, fill_rgb=TEAL)

# constellation dot grid (decorative)
for row in range(5):
    for col in range(6):
        cx = Inches(0.5 + col * 0.7)
        cy = Inches(0.6 + row * 1.1)
        dot = s.shapes.add_shape(9, cx, cy, Inches(0.08), Inches(0.08))  # oval
        dot.fill.solid(); dot.fill.fore_color.rgb = RGBColor(0xFF,0xFF,0xFF)
        dot.fill.fore_color.theme_color  # ignore; just set alpha via xml workaround
        dot.line.fill.background()

# connecting lines between some dots (simple star pattern)
def add_line(slide, x1, y1, x2, y2, color=WHITE):
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width     = Pt(0.75)

add_line(s, Inches(0.78), Inches(0.64), Inches(1.48), Inches(1.74))
add_line(s, Inches(1.48), Inches(1.74), Inches(2.18), Inches(0.64))
add_line(s, Inches(2.18), Inches(0.64), Inches(2.88), Inches(1.74))
add_line(s, Inches(0.78), Inches(2.84), Inches(1.48), Inches(1.74))
add_line(s, Inches(2.88), Inches(1.74), Inches(2.18), Inches(2.84))

# wordmark
add_text(s, "MindCraft", Inches(5.5), Inches(1.8), Inches(7.5), Inches(1.2),
         font_size=60, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s, "Redefining Student Success",
         Inches(5.5), Inches(2.9), Inches(7.2), Inches(0.6),
         font_size=22, color=GOLD, align=PP_ALIGN.LEFT)

add_rect(s, Inches(5.5), Inches(3.6), Inches(4.5), Inches(0.04), fill_rgb=LIGHT)

add_text(s, "BETA Accelerator — Company Introduction",
         Inches(5.5), Inches(3.8), Inches(7.2), Inches(0.5),
         font_size=14, color=LIGHT, italic=True, align=PP_ALIGN.LEFT)

add_text(s, "Akshat Koirala & Kaif  |  2026",
         Inches(5.5), Inches(6.5), Inches(7.2), Inches(0.5),
         font_size=13, color=LIGHT, align=PP_ALIGN.LEFT)

accent_bar(s, color=GOLD)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 – THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s, color=TEAL)
slide_number(s, 2)

add_text(s, "THE PROBLEM", Inches(0.7), Inches(0.4), Inches(5), Inches(0.5),
         font_size=12, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
add_text(s, "Schools teach in straight lines.\nLearning doesn't work that way.",
         Inches(0.7), Inches(0.9), Inches(7.5), Inches(1.4),
         font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# story card
add_rect(s, Inches(0.7), Inches(2.5), Inches(7.6), Inches(3.5),
         fill_rgb=RGBColor(0x14, 0x2A, 0x40))
add_text(s, "Meet Seb.", Inches(1.0), Inches(2.7), Inches(7), Inches(0.5),
         font_size=16, bold=True, color=GOLD)
story = (
    "Seb understands derivatives in calculus — but still gets stuck\n"
    "on the algebra underneath it. On paper, he looked fine.\n"
    "Teachers said he was doing well.\n\n"
    "But his parents saw his confidence quietly fading as the gaps\n"
    "beneath the surface started to show."
)
add_text(s, story, Inches(1.0), Inches(3.2), Inches(7.2), Inches(2.5),
         font_size=15, color=LIGHT)

# right stat block
add_rect(s, Inches(9.2), Inches(1.0), Inches(3.7), Inches(5.5),
         fill_rgb=RGBColor(0x1A, 0x35, 0x50))
add_text(s, "The Gap Is Real", Inches(9.4), Inches(1.2), Inches(3.3), Inches(0.5),
         font_size=14, bold=True, color=GOLD)
stats = [
    ("70 %", "of students have\nunaddressed gaps\nfrom COVID-era\nlearning loss"),
    ("1 in 3", "students can't\nafford premium\ntutoring support"),
    ("< 10 %", "of tutoring is\nstructured around\nhow understanding\nactually builds"),
]
for i, (num, label) in enumerate(stats):
    yoff = Inches(1.85 + i * 1.55)
    add_text(s, num, Inches(9.4), yoff, Inches(1.3), Inches(0.6),
             font_size=28, bold=True, color=TEAL)
    add_text(s, label, Inches(10.75), yoff, Inches(2.0), Inches(0.7),
             font_size=11, color=LIGHT)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 – OUR SOLUTION
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s, color=GOLD)
slide_number(s, 3)

add_text(s, "OUR SOLUTION", Inches(0.7), Inches(0.4), Inches(5), Inches(0.5),
         font_size=12, bold=True, color=TEAL)
add_text(s, "Three pillars. One coherent learning experience.",
         Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
         font_size=30, bold=True, color=WHITE)

pillars = [
    (TEAL,  "01", "Adaptive\nTutoring",
     "One-on-one sessions that adapt in real time to where the student actually is — not where the curriculum says they should be."),
    (GOLD,  "02", "Targeted\nPractice",
     "AI-generated problem sets that target the exact gap, not the whole topic. Practice meets the student, not the other way around."),
    (RGBColor(0x7B,0x5E,0xA7), "03", "Student\nConstellations",
     "Visual maps of what a student understands, where they struggle, and what they're ready to learn next — updated after every session."),
]
for i, (col, num, title, body) in enumerate(pillars):
    x = Inches(0.6 + i * 4.2)
    add_rect(s, x, Inches(2.0), Inches(3.85), Inches(4.8),
             fill_rgb=RGBColor(0x14, 0x28, 0x3E))
    add_rect(s, x, Inches(2.0), Inches(0.5), Inches(0.12), fill_rgb=col)
    add_text(s, num, x + Inches(0.2), Inches(2.2), Inches(1), Inches(0.5),
             font_size=28, bold=True, color=col)
    add_text(s, title, x + Inches(0.2), Inches(2.85), Inches(3.4), Inches(0.75),
             font_size=20, bold=True, color=WHITE)
    add_text(s, body, x + Inches(0.2), Inches(3.65), Inches(3.4), Inches(2.8),
             font_size=13, color=LIGHT)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 – WHO WE SERVE
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s, color=TEAL)
slide_number(s, 4)

add_text(s, "WHO WE SERVE", Inches(0.7), Inches(0.4), Inches(5), Inches(0.5),
         font_size=12, bold=True, color=GOLD)
add_text(s, "We solve for the full learning ecosystem.",
         Inches(0.7), Inches(0.85), Inches(11), Inches(0.7),
         font_size=30, bold=True, color=WHITE)

segments = [
    (GOLD,  "Parents",
     "Families who want high-quality academic support without paying $80–120/hr for "
     "traditional tutoring. They need transparency, results, and reassurance."),
    (TEAL,  "Students",
     "Learners who need more than a worksheet. They need structure, clarity, and "
     "a tutor who meets them where they actually are — not where the syllabus expects."),
    (RGBColor(0x7B,0x5E,0xA7), "Tutors",
     "Educators who want to build credibility, grow a client base, and develop "
     "their careers within a structured, supportive platform."),
]
for i, (col, title, body) in enumerate(segments):
    x = Inches(0.6 + i * 4.2)
    # circle icon area
    circ = s.shapes.add_shape(9, x + Inches(1.4), Inches(1.9), Inches(1.0), Inches(1.0))
    circ.fill.solid(); circ.fill.fore_color.rgb = col
    circ.line.fill.background()
    add_text(s, title, x, Inches(3.1), Inches(3.85), Inches(0.55),
             font_size=20, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_rect(s, x + Inches(1.5), Inches(3.7), Inches(0.85), Inches(0.06), fill_rgb=col)
    add_text(s, body, x + Inches(0.1), Inches(3.95), Inches(3.65), Inches(2.5),
             font_size=13, color=LIGHT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 – COMPETITION & DIFFERENTIATION
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s, color=GOLD)
slide_number(s, 5)

add_text(s, "COMPETITION & DIFFERENTIATION", Inches(0.7), Inches(0.4),
         Inches(8), Inches(0.5), font_size=12, bold=True, color=TEAL)
add_text(s, "Crowded market. Different model.",
         Inches(0.7), Inches(0.85), Inches(9), Inches(0.7),
         font_size=30, bold=True, color=WHITE)

# comparison table
headers = ["", "Khan Academy", "Varsity Tutors", "Wyzant", "MindCraft"]
rows = [
    ("Live tutoring",       "✗", "✓", "✓", "✓"),
    ("Structured practice", "✓", "✗", "✗", "✓"),
    ("Knowledge mapping",   "✗", "✗", "✗", "✓"),
    ("Affordable pricing",  "✓", "✗", "✗", "✓"),
    ("Tutor career tools",  "✗", "✓", "Limited", "✓"),
]
col_colors = [NAVY, NAVY, NAVY, NAVY, TEAL]
col_xs = [Inches(0.5), Inches(2.9), Inches(5.0), Inches(7.2), Inches(9.4)]
col_ws = [Inches(2.3), Inches(2.0), Inches(2.0), Inches(2.0), Inches(3.5)]

# header row
for j, (hdr, cx, cw, cc) in enumerate(zip(headers, col_xs, col_ws, col_colors)):
    bg_c = TEAL if j == 4 else RGBColor(0x14, 0x28, 0x3E)
    add_rect(s, cx, Inches(1.75), cw, Inches(0.45), fill_rgb=bg_c)
    add_text(s, hdr, cx + Inches(0.08), Inches(1.80), cw - Inches(0.15), Inches(0.38),
             font_size=13, bold=True,
             color=NAVY if j == 4 else GOLD,
             align=PP_ALIGN.CENTER)

for i, row in enumerate(rows):
    row_y = Inches(2.25 + i * 0.72)
    for j, (val, cx, cw) in enumerate(zip(row, col_xs, col_ws)):
        bg_c = RGBColor(0x18, 0x32, 0x4A) if i % 2 == 0 else RGBColor(0x14, 0x28, 0x3E)
        if j == 4: bg_c = RGBColor(0x18, 0x40, 0x50)
        add_rect(s, cx, row_y, cw, Inches(0.65), fill_rgb=bg_c)
        txt_color = GOLD if j == 4 else (WHITE if val == "✓" else RGBColor(0x99,0xAA,0xBB))
        add_text(s, val, cx + Inches(0.08), row_y + Inches(0.12),
                 cw - Inches(0.15), Inches(0.45),
                 font_size=13 if j == 0 else 14,
                 bold=(j == 4),
                 color=txt_color,
                 align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 – BUSINESS MODEL (AFFORDABILITY)
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s, color=TEAL)
slide_number(s, 6)

add_text(s, "BUSINESS MODEL", Inches(0.7), Inches(0.4), Inches(5), Inches(0.5),
         font_size=12, bold=True, color=GOLD)
add_text(s, "Premium support. Accessible price.",
         Inches(0.7), Inches(0.85), Inches(9), Inches(0.7),
         font_size=30, bold=True, color=WHITE)

model_cards = [
    (TEAL,  "Subscription\nPlans",
     "Monthly bundles covering live sessions + practice + notes.\n"
     "Priced at a fraction of traditional tutoring."),
    (GOLD,  "Tutor Revenue\nShare",
     "Tutors earn per session and grow their client base\n"
     "through the platform — aligned incentives."),
    (RGBColor(0x7B,0x5E,0xA7), "Institution\nPartnerships",
     "Schools & districts license MindCraft to close\n"
     "learning gaps at scale — B2B2C path."),
]
for i, (col, title, body) in enumerate(model_cards):
    x = Inches(0.6 + i * 4.2)
    add_rect(s, x, Inches(1.9), Inches(3.85), Inches(4.5),
             fill_rgb=RGBColor(0x14, 0x28, 0x3E))
    add_rect(s, x, Inches(1.9), Inches(3.85), Inches(0.12), fill_rgb=col)
    add_text(s, title, x + Inches(0.2), Inches(2.15), Inches(3.4), Inches(0.8),
             font_size=19, bold=True, color=col)
    add_text(s, body, x + Inches(0.2), Inches(3.1), Inches(3.4), Inches(2.8),
             font_size=13, color=LIGHT)

add_text(s,
         "Traditional premium tutoring: $80–$120 / hr\n"
         "MindCraft target: $30–50 / hr equivalent  ✓  No quality compromise",
         Inches(0.7), Inches(6.4), Inches(12), Inches(0.75),
         font_size=13, color=GOLD, italic=True)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 – TRACTION & VALIDATION
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s, color=GOLD)
slide_number(s, 7)

add_text(s, "TRACTION & VALIDATION", Inches(0.7), Inches(0.4), Inches(7), Inches(0.5),
         font_size=12, bold=True, color=TEAL)
add_text(s, "Early signal. Real conversations.",
         Inches(0.7), Inches(0.85), Inches(9), Inches(0.7),
         font_size=30, bold=True, color=WHITE)

milestones = [
    (TEAL,  "Product",      "Working MVP: adaptive tutoring, practice generation,\nstudent constellation dashboard — deployed."),
    (GOLD,  "Mentorship",   "Advising from Allan Martinez (Director, Macalester\nEntrepreneurship & Idea Lab) and Erik Halaas."),
    (RGBColor(0x7B,0x5E,0xA7), "Discovery", "Direct interviews with parents and tutors —\npractical and emotional needs documented."),
    (TEAL,  "Learning",     "Core insight: strong companies are built by listening\nclosely, not by chasing the coolest features."),
]
for i, (col, label, text) in enumerate(milestones):
    row = i // 2
    col_i = i % 2
    x = Inches(0.6 + col_i * 6.4)
    y = Inches(2.0 + row * 2.15)
    add_rect(s, x, y, Inches(6.0), Inches(1.85),
             fill_rgb=RGBColor(0x14, 0x28, 0x3E))
    add_rect(s, x, y, Inches(0.12), Inches(1.85), fill_rgb=col)
    add_text(s, label, x + Inches(0.28), y + Inches(0.15), Inches(5.5), Inches(0.4),
             font_size=14, bold=True, color=col)
    add_text(s, text, x + Inches(0.28), y + Inches(0.6), Inches(5.5), Inches(1.1),
             font_size=13, color=LIGHT)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8 – FOUNDER INTRODUCTIONS
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s, color=TEAL)
slide_number(s, 8)

add_text(s, "THE FOUNDERS", Inches(0.7), Inches(0.4), Inches(5), Inches(0.5),
         font_size=12, bold=True, color=GOLD)
add_text(s, "Built from lived experience in education.",
         Inches(0.7), Inches(0.85), Inches(11), Inches(0.7),
         font_size=30, bold=True, color=WHITE)

founders = [
    ("Akshat Koirala", "Co-Founder",
     [
         "Macalester College — Economics & Statistics",
         "Product & engineering lead for MindCraft",
         "Built the MVP: tutoring flow, practice engine,\nconstellation dashboard",
         "Motivated by watching talented students fall\nbehind not from lack of effort, but lack of\nthe right support structure",
     ]),
    ("Kaif", "Co-Founder",
     [
         "Macalester College",
         "Strategy, operations & customer discovery lead",
         "Conducted in-depth interviews with parents\nand tutors across income brackets",
         "Passionate about closing the opportunity gap\nin academic support for underserved families",
     ]),
]
for i, (name, role, bullets) in enumerate(founders):
    x = Inches(0.5 + i * 6.5)
    # avatar placeholder
    av = s.shapes.add_shape(9, x + Inches(0.3), Inches(1.85), Inches(1.5), Inches(1.5))
    av.fill.solid(); av.fill.fore_color.rgb = TEAL if i == 0 else GOLD
    av.line.fill.background()
    add_text(s, name[0], x + Inches(0.7), Inches(2.1), Inches(0.9), Inches(0.9),
             font_size=36, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, name, x + Inches(2.0), Inches(1.9), Inches(4.0), Inches(0.55),
             font_size=20, bold=True, color=WHITE)
    add_text(s, role, x + Inches(2.0), Inches(2.45), Inches(4.0), Inches(0.4),
             font_size=13, color=GOLD, italic=True)
    for j, b in enumerate(bullets):
        add_text(s, f"• {b}", x + Inches(0.3), Inches(3.55 + j * 0.72),
                 Inches(5.9), Inches(0.65),
                 font_size=12, color=LIGHT)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 9 – WHY BETA
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s, color=GOLD)
slide_number(s, 9)

add_text(s, "WHY BETA?", Inches(0.7), Inches(0.4), Inches(5), Inches(0.5),
         font_size=12, bold=True, color=TEAL)
add_text(s, "We are eager to learn,\neager to build.",
         Inches(0.7), Inches(0.85), Inches(7), Inches(1.3),
         font_size=34, bold=True, color=WHITE)

needs = [
    (TEAL,  "Design Thinking\n& Customer Feedback",
     "We want to keep shaping MindCraft around real user needs — BETA's "
     "curriculum gives us the tools to do that rigorously."),
    (GOLD,  "Financial Models\n& Valuation",
     "We don't yet have deep experience thinking through capital and "
     "long-term sustainability. We need that foundation."),
    (RGBColor(0x7B,0x5E,0xA7), "Mentor Network",
     "Pressure-testing our thinking with experienced operators "
     "will sharpen us in ways we cannot replicate on our own."),
]
for i, (col, title, body) in enumerate(needs):
    x = Inches(0.6 + i * 4.2)
    add_rect(s, x, Inches(2.5), Inches(3.85), Inches(4.0),
             fill_rgb=RGBColor(0x14, 0x28, 0x3E))
    add_rect(s, x, Inches(2.5), Inches(3.85), Inches(0.1), fill_rgb=col)
    add_text(s, title, x + Inches(0.2), Inches(2.75), Inches(3.4), Inches(0.8),
             font_size=16, bold=True, color=col)
    add_text(s, body, x + Inches(0.2), Inches(3.65), Inches(3.4), Inches(2.5),
             font_size=13, color=LIGHT)

add_text(s,
         '"Strong companies are not built by chasing the coolest features,\n'
         'but by listening closely, learning quickly, and being willing to revise your thinking."',
         Inches(0.7), Inches(6.5), Inches(12), Inches(0.85),
         font_size=13, color=GOLD, italic=True)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 10 – CLOSE / CALL TO ACTION
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

# full-width teal strip top third
add_rect(s, 0, 0, SLIDE_W, Inches(2.6), fill_rgb=TEAL)

add_text(s, "Our goal is not simply to help students keep up —",
         Inches(0.8), Inches(0.3), Inches(12), Inches(0.7),
         font_size=24, bold=True, color=NAVY)
add_text(s, "but to help more of them genuinely thrive.",
         Inches(0.8), Inches(0.95), Inches(12), Inches(0.7),
         font_size=24, bold=True, color=NAVY)
add_text(s, "Especially those who may not otherwise have access to high-quality academic support.",
         Inches(0.8), Inches(1.6), Inches(12), Inches(0.6),
         font_size=16, color=NAVY, italic=True)

add_text(s, "MindCraft", Inches(0.8), Inches(3.2), Inches(6), Inches(0.9),
         font_size=52, bold=True, color=WHITE)
add_text(s, "Let's build it together.", Inches(0.8), Inches(4.1), Inches(8), Inches(0.6),
         font_size=22, color=GOLD)

add_rect(s, Inches(0.8), Inches(4.85), Inches(4), Inches(0.06), fill_rgb=TEAL)

add_text(s, "Akshat Koirala   |   akoirala@macalester.edu",
         Inches(0.8), Inches(5.1), Inches(9), Inches(0.45),
         font_size=14, color=LIGHT)
add_text(s, "mindcraft.study",
         Inches(0.8), Inches(5.55), Inches(4), Inches(0.4),
         font_size=14, color=TEAL)

# large decorative "M"
add_text(s, "M", Inches(9.5), Inches(2.5), Inches(3.5), Inches(4.5),
         font_size=260, bold=True, color=RGBColor(0x1A, 0x35, 0x50), align=PP_ALIGN.CENTER)

accent_bar(s, color=GOLD)


# ── save ───────────────────────────────────────────────────────────────────
out = "/Users/akoirala/Desktop/Business Ideas/mindcraft-site/MindCraft_BETA_Deck.pptx"
prs.save(out)
print(f"Saved → {out}")
